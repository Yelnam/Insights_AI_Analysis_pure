# -*- coding: utf-8 -*-
"""
Created on Mon Nov 21 13:23:18 2022

@author: rober
"""
import pandas as pd

from utils.dicts_data_to_ppt import dict_low_numbers
from utils.func_data import create_df_author_sources
from utils.func_generate_headlines import generate_headlines

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


df_temp = pd.DataFrame()

# %% temporary data values ----------------------------------------------------

# temporary values in order to avoid errors in functions below
# these will be overwritten by actual datapoints in main script
year = 'year'
company = 'company'
company_short = 'company_short'
top_author_by_vol = 'top_author_by_vol'
top_author_by_vol_surname = 'top_author_by_vol_surname'
top_author_by_vol_source = 'top_author_by_vol_source'
total_vol_text = 'total_vol_text'
max_m_vol_text_long = 'max_m_vol_text_long'
max_vol_month_text = 'max_vol_month_text'
total_reach_text_long = 'total_reach_text_long'
max_m_reach_text_long = 'max_m_reach_text_long'
max_reach_month_text = 'max_reach_month_text'
top_media_type = 'top_media_type'
top_media_type_vol_text = 'top_media_type_vol_text'
top_source_by_vol = 'top_source_by_vol'
top_source_vol_text = 'top_source_vol_text'
top_source_by_reach = 'top_source_by_reach'
top_source_reach_text_long = 'top_source_reach_text_long'
top_author_vol = 'top_author_vol'
top_author_reach_text_long = 'top_author_reach_text_long'
top_author_reach_pct = 'top_author_reach_pct'
top_author_by_reach = 'top_author_by_reach'
top_author_by_reach_surname = 'top_author_by_reach_surname'
top_author_by_reach_source = 'top_author_by_reach_source'
top_author_reach = 'top_author_reach'
top_source_top_author = 'top_source_top_author'
top_source_top_author_vol = 'top_source_top_author_vol'
top_source_reach_text_long = 'top_source_reach_text_long'
second_source_by_reach = 'second_source_by_reach'
top_source_reach_pct_vs_2 = 'top_source_reach_pct_vs_2'
top_source_reach_pct_vs_vol = 'top_source_reach_pct_vs_vol'
top_source_months_at_top = 'top_source_months_at_top'

# %% add text general ---------------------------------------------------------

def add_text(slide, 
             text = 'Default text', 
             left = Inches(1), 
             top = Inches(1), 
             width = Inches(1), 
             height = Inches(1), 
             font = 'Arial',
             size = Pt(18), 
             color = RGBColor(000, 000, 000),
             wrap = True, 
             align = PP_ALIGN.LEFT):
    
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    textbox = textbox.paragraphs[0]
    
    # add text
    textbox.text = text
    
    # format text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    
    # set alignment
    textbox.alignment = align
    
# %% add EOY title ------------------------------------------------------------
    
def add_text_EOY2022_title(slide, 
             text = 'Default text', 
             left = Inches(0.41), 
             top = Inches(0.38), 
             width = Inches(3.5), 
             height = Inches(0.51), 
             font = 'Bahnschrift Light', 
             size = 24, 
             color = RGBColor(000, 000, 000),
             wrap = False,
             align = PP_ALIGN.LEFT):
    
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    textbox = textbox.paragraphs[0]
    
    # add text
    textbox.text = text
    
    # format text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    
    # set alignment
    textbox.alignment = align
    
# %% add EOY header -----------------------------------------------------------    
        
def add_text_EOY2022_header(slide, 
             text = 'Default text', 
             left = Inches(4), 
             top = Inches(0.5), 
             width = Inches(2.75), 
             height = Inches(0.34), 
             font = 'Bahnschrift Light', 
             size = 14, 
             color = RGBColor(000, 000, 000),
             wrap = False,
             align = PP_ALIGN.LEFT):
    
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    textbox = textbox.paragraphs[0]
    
    # add text
    textbox.text = text
    
    # format text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    
    # set alignment
    textbox.alignment = align

# %% add EOY topline figures --------------------------------------------------

def add_text_EOY2022_topline(slide, 
             text = 'Default text', 
             left = Inches(4), 
             top = Inches(0.5), 
             width = Inches(1.3), 
             height = Inches(0.71), 
             font = 'Bahnschrift SemiBold', 
             size = 36, 
             color = RGBColor(248, 166, 98),
             wrap = False,
             align = PP_ALIGN.CENTER):
    
    # creating a text box
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    textbox = textbox.paragraphs[0]
    
    # add text
    textbox.text = text
    
    # format text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    
    # set alignment
    textbox.alignment = align

# %% add EOY body text --------------------------------------------------------
    
def add_text_EOY2022_body(slide, 
            text = 'Default text', 
            left = Inches(1.5), 
            top = Inches(1.67), 
            width = Inches(5.54), 
            height = Inches(1.2), 
            font = 'Bahnschrift Light', 
            size = 16, 
            color = RGBColor(000, 000, 000),
            wrap = True,
            align = PP_ALIGN.LEFT):
        
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    textbox = textbox.paragraphs[0]
    
    # add text
    textbox.text = text
    
    # format text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    
    # set alignment
    textbox.alignment = align
    
# %% add EOY headlines --------------------------------------------------------
    
def add_text_EOY2022_headlines(slide = 'slide', dataframe = df_temp,
        co= 'a', co_hl = 'b', co_alt = 'c',
        left = Inches(9.88), top = Inches(3.59), 
        width = Inches(3.01), height = Inches(3.59), 
        font = 'Bierstadt Display', size = Pt(11), 
        color = RGBColor(89, 89, 89), wrap = True, align = PP_ALIGN.CENTER):
    
    (hl1, hl1_source, hl2, hl2_source, 
    hl3, hl3_source, hl4, hl4_source) = generate_headlines(dataframe, co, co_hl, co_alt)
    
    txt_top_hls = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_top_hls.word_wrap = wrap
    txt_top_hls = txt_top_hls.paragraphs[0]
    txt_top_hls.alignment = align
    txt_top_hls.font.name = font
    txt_top_hls.font.color.rgb = color
    txt_top_hls.font.size = size

    run1hl, run1hl.text = txt_top_hls.add_run(), f'{hl1_source}\n'
    run2hl, run2hl.text = txt_top_hls.add_run(), f'{hl1}\n\n'
    run3hl, run3hl.text = txt_top_hls.add_run(), f'{hl2_source}\n'
    run4hl, run4hl.text = txt_top_hls.add_run(), f'{hl2}\n\n'
    run5hl, run5hl.text = txt_top_hls.add_run(), f'{hl3_source}\n'
    run6hl, run6hl.text = txt_top_hls.add_run(), f'{hl3}\n\n'
    run7hl, run7hl.text = txt_top_hls.add_run(), f'{hl4_source}\n'
    run8hl, run8hl.text = txt_top_hls.add_run(), f'{hl4}\n\n'

    # format text
    run1hl.font.bold = run3hl.font.bold = run5hl.font.bold = run7hl.font.bold = True
    run1hl.font.italic = run3hl.font.italic = run5hl.font.italic = run7hl.font.italic = True
    run1hl.font.color.rgb = run3hl.font.color.rgb = run5hl.font.color.rgb = run7hl.font.color.rgb = RGBColor(55,96,146)
    run1hl.font.size = run3hl.font.size = run5hl.font.size = run7hl.font.size = Pt(16)

# %% add EOY top source text --------------------------------------------------

def add_text_EOY2022_top_source(slide, text = 'Default text', 
        top_source_by_vol = 'top_source_by_vol',
        top_source_vol_text = 'top_source_vol_text',
        left = Inches(6.75), top = Inches(3.43), 
        width = Inches(2.85), height = Inches(0.59), 
        font = 'Bahnschrift Light', size = 16, 
        color = RGBColor(000, 000, 000), wrap = True, align = PP_ALIGN.LEFT):
    
    # EXEC SUMMARY - TOP SOURCE
    txt_top_source = (slide.shapes.add_textbox(left, top, width, height)).text_frame
                      
    txt_top_source.word_wrap = True
    txt_top_source = txt_top_source.paragraphs[0]
    txt_top_source.alignment = PP_ALIGN.CENTER

    run1s = txt_top_source.add_run()
    run1s.text = f'{top_source_by_vol}\n'
    run2s = txt_top_source.add_run()
    run2s.text = f'{top_source_vol_text} articles'

    run1s.font.italic = True
    run1s.font.bold = True
    run1s.font.name = 'Bierstadt Display'
    run1s.font.color.rgb = RGBColor(0,176,80)
    run1s.font.size = Pt(18) if len(top_source_by_vol) <= 20 else Pt(14) # sets smaller text size if long source name
    run2s.font.name = 'Bierstadt Display'
    run2s.font.color.rgb = RGBColor(89,89,89)
    run2s.font.size = Pt(11)

# %% add EOY top author text --------------------------------------------------

def add_text_EOY2022_top_author(slide, 
        top_author_by_vol = 'top_author_by_vol',
        top_author_vol_text = 'top_author_vol_text',
        top_author_by_vol_source = 'top_author_by_vol_source',
        left = Inches(6.75), top = Inches(4.36), 
        width = Inches(2.85), height = Inches(0.77), 
        font = 'Bahnschrift Light', size = 16, 
        color = RGBColor(000, 000, 000), wrap = True, align = PP_ALIGN.LEFT):
    
    # EXEC SUMMARY - TOP AUTHOR
    txt_top_author = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_top_author.word_wrap = True
    txt_top_author = txt_top_author.paragraphs[0]
    txt_top_author.alignment = PP_ALIGN.CENTER
    txt_top_author.font.name = 'Bierstadt Display'

    run1a = txt_top_author.add_run()
    run1a.text = f'{top_author_by_vol}\n'
    run2a = txt_top_author.add_run()
    run2a.text = f'{top_author_by_vol_source}\n'
    run3a = txt_top_author.add_run()
    run3a.text = f'{top_author_vol_text} articles'

    # format text
    run1a.font.bold = True
    run1a.font.color.rgb = RGBColor(0,176,80)
    run1a.font.size = Pt(18) if len(top_author_by_vol) <= 20 else Pt(14) # sets smaller text size if long author name
    run2a.font.italic = True
    run2a.font.color.rgb = RGBColor(0,0,0)
    run2a.font.size = Pt(11)
    run3a.font.color.rgb = RGBColor(89,89,89)
    run3a.font.size = Pt(11)

# %% add EOY chart decoration -------------------------------------------------

def add_text_EOY2022_chart_decor(slide, 
            text = 'Default text', 
            left = Inches(3.99), 
            top = Inches(0.81), 
            width = Inches(0.25), 
            height = Inches(1.2), 
            font = 'Calibri', 
            size = 9, 
            color = RGBColor(166, 166, 166),
            wrap = True,
            align = PP_ALIGN.RIGHT):
    
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    # add paragraph
    textbox = textbox.paragraphs[0]
    
    # format
    textbox.text = text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    textbox.alignment = align
    
# %% add EOY monthly breakdown side headers ----------------------------------    
    
def add_text_EOY2022_m_side_headers(slide, 
            text = 'Default text', 
            left = Inches(0.49), 
            top = Inches(2.21), 
            width = Inches(1.35), 
            height = Inches(0.3), 
            font = 'Bahnschrift Light', 
            size = 12, 
            color = RGBColor(122, 122, 122),
            wrap = True,
            align = PP_ALIGN.RIGHT):
    
    # creating a text box and adding paragraph
    textbox = (slide.shapes.add_textbox
                      (left, top, width, height)).text_frame
    
    # set word wrap (NOT SURE WHY, BUT MUST BE DONE BEFORE ADDING PARAGRAPH)
    textbox.word_wrap = wrap
    
    # add paragraph
    textbox = textbox.paragraphs[0]
    
    # format
    textbox.text = text
    textbox.font.name = font
    textbox.font.size = Pt(size)
    textbox.font.color.rgb = color
    textbox.alignment = align    

# %% add EOY executive summmary -----------------------------------------------

def add_text_EOY2022_summary(slide, 
            company = company, # variables set as themselves in text form are placeholders
            company_short = company_short, # these are to be replaced with the actual variables
            year = year,
            total_vol_text = total_vol_text,
            max_m_vol_text_long = max_m_vol_text_long,
            max_vol_month_text = max_vol_month_text,
            total_reach_text_long = total_reach_text_long,
            max_m_reach_text_long = max_m_reach_text_long,
            max_reach_month_text = max_reach_month_text,
            top_media_type = top_media_type,
            top_media_type_vol_text = top_media_type_vol_text,
            top_source_by_vol = top_source_by_vol,
            top_source_vol_text = top_source_vol_text,
            top_source_by_reach = top_source_by_reach,
            top_source_reach_text_long = top_source_reach_text_long,
            left = Inches(0.6), 
            top = Inches(2.43), 
            width = Inches(3.2), 
            height = Inches(4.14), 
            font = 'Bahnschrift Light', 
            size = 12,
            wrap = True,
            align = PP_ALIGN.LEFT):
     
    txt_exec_summ = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_exec_summ.word_wrap = wrap
    txt_exec_summ = txt_exec_summ.paragraphs[0]
    txt_exec_summ.font.name = font
    txt_exec_summ.font.size = Pt(size)

    # creating runs (sections of text that can be individually formatted)
    run1 = txt_exec_summ.add_run()
    run1.text = f'● {company} generated {total_vol_text} mentions across {year}, reaching a potential audience of {total_reach_text_long}.\n\n'

    run2 = txt_exec_summ.add_run()
    run2.text = f'● Monthly volume peaked at {max_m_vol_text_long} articles in {max_vol_month_text}, with {year}\'s largest monthly audience of {max_m_reach_text_long} reached in '

    run3 = txt_exec_summ.add_run()
    run3.text = 'the same month.\n\n' if max_vol_month_text == max_reach_month_text else f'{max_reach_month_text}.\n\n'

    run4 = txt_exec_summ.add_run()
    run4.text = f'● Mentions appeared most often in {top_media_type} media, which accounted for {top_media_type_vol_text} items, while the leading publication by volume was '

    run5 = txt_exec_summ.add_run()
    run5.text = f'{top_source_by_vol}, '

    run6 = txt_exec_summ.add_run()
    run6.text = f'which mentioned {company_short} {top_source_vol_text} times this year.\n\n'

    run7 = txt_exec_summ.add_run()
    run8 = txt_exec_summ.add_run()
    if top_source_by_vol == top_source_by_reach:
        run7.text = f'● {top_source_by_reach}' 
        run8.text = ' was also the top source by reach, '
    else:
        run7.text = f'● {top_source_by_reach}'
        run8.text = ' was the top source by reach, '

    run9 = txt_exec_summ.add_run()
    run9.text = f'with a potential audience of {top_source_reach_text_long} exposed to {company_short} mentions through its coverage.'

    # formatting above runs
    run5.font.italic = True
    run7.font.italic = True

# %% add EOY text, top author by vol ------------------------------------------

def add_text_author_by_vol(slide, 
                           dataframe = df_temp,
            company = company, # variables set as themselves in text form are placeholders
            company_short = company_short, # these are to be replaced with the actual variables
            top_author_by_vol = top_author_by_vol,
            top_author_by_vol_surname = top_author_by_vol_surname,
            top_author_by_vol_source = top_author_by_vol_source,
            top_author_vol = top_author_vol,
            left = Inches(1.5), 
            top = Inches(1.67), 
            width = Inches(5.54), 
            height = Inches(1.2), 
            font = 'Bahnschrift Light', 
            size = 16,
            wrap = True,
            align = PP_ALIGN.LEFT):
     
    txt_auth_vol = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_auth_vol.word_wrap = wrap
    txt_auth_vol = txt_auth_vol.paragraphs[0]
    txt_auth_vol.font.name = font
    txt_auth_vol.font.size = Pt(size)

    # creating runs (sections of text that can be individually formatted)
    run1 = txt_auth_vol.add_run()
    run1.text = f'{top_author_by_vol} wrote about {company} more often than any other journalist. '

    run2 = txt_auth_vol.add_run()
    run2.text = ('Writing exclusively for ' 
                 if len(create_df_author_sources(dataframe, top_author_by_vol)) == 1 
                 else 'Writing primarily for ')

    run3 = txt_auth_vol.add_run()
    run3.text = f'{top_author_by_vol_source}'
    
    run4 = txt_auth_vol.add_run()
    run4.text = f', {top_author_by_vol_surname} penned {top_author_vol} items which mentioned {company_short} by name.'

    # formatting above runs
    run3.font.italic = True

# %% add EOY text, top author by reach ----------------------------------------

def add_text_author_by_reach(slide, 
                           dataframe = df_temp,
            company = company, # variables set as themselves in text form are placeholders
            company_short = company_short, # these are to be replaced with the actual variables
            top_author_by_vol = top_author_by_vol,
            top_author_by_vol_surname = top_author_by_vol_surname,
            top_author_vol = top_author_vol,
            top_author_reach_text_long = top_author_reach_text_long,
            top_author_reach_pct = top_author_reach_pct,
            top_author_by_reach = top_author_by_reach,
            top_author_by_reach_surname = top_author_by_reach_surname,
            top_author_by_reach_source = top_author_by_reach_source,
            top_author_reach = top_author_reach,
            left = Inches(1.5), 
            top = Inches(4.73), 
            width = Inches(5.54), 
            height = Inches(1.2), 
            font = 'Bahnschrift Light', 
            size = 16,
            wrap = True,
            align = PP_ALIGN.LEFT):
     
    txt_auth_reach = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_auth_reach.word_wrap = wrap
    txt_auth_reach = txt_auth_reach.paragraphs[0]
    txt_auth_reach.font.name = font
    txt_auth_reach.font.size = Pt(size)

    # creating runs (sections of text that can be individually formatted)
    if top_author_by_vol == top_author_by_reach:
        run1 = txt_auth_reach.add_run()
        run1.text = f'{top_author_by_vol_surname} was also the widest-reaching journalist. His {top_author_vol} articles reached a total potential readership of {top_author_reach_text_long}, or {top_author_reach_pct}% of the total reach for {company_short} mentions.'

    else:
        if len(create_df_author_sources(dataframe, top_author_by_reach)) == 1:
            run1 = txt_auth_reach.add_run()
            run1.text = f'While {top_author_by_vol_surname} was the most prominent journalist by volume, {top_author_by_reach} was the widest-reaching. Writing exclusively for '
            
            run2 = txt_auth_reach.add_run()
            run2.text =  f'{top_author_by_reach_source}'
            
            run3 = txt_auth_reach.add_run()
            run3.text =  f', {top_author_by_reach_surname} reached {top_author_reach_text_long} readers, or {top_author_reach_pct}% of the total reach for {company_short} mentions.'

            # formatting above runs
            run2.font.italic = True
            
        else:
            run1 = txt_auth_reach.add_run()
            run1.text =  f'While {top_author_by_vol_surname} was the most prominent journalist by volume, {top_author_by_reach} was the widest-reaching. Writing primarily for '
            
            run2 = txt_auth_reach.add_run()
            run2.text =  f'{top_author_by_reach_source}'
            
            run3 = txt_auth_reach.add_run()
            run3.text =  f', {top_author_by_reach_surname} reached {top_author_reach_text_long} readers, or {top_author_reach_pct}% of the total reach for {company_short} mentions.'
    
            # formatting above runs
            run2.font.italic = True

# %% add EOY text, top source by vol ------------------------------------------

def add_text_source_by_vol(slide, 
            dataframe = df_temp,
            year = year,
            company = company, # variables set as themselves in text form are placeholders
            company_short = company_short, # these are to be replaced with the actual variables
            top_source_by_vol = top_source_by_vol,
            top_source_top_author = top_source_top_author,
            top_source_top_author_vol = top_source_top_author_vol,
            left = Inches(1.5), 
            top = Inches(1.67), 
            width = Inches(5.54), 
            height = Inches(1.2), 
            font = 'Bahnschrift Light', 
            size = 16,
            wrap = True,
            align = PP_ALIGN.LEFT):
    
    txt_source_vol = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_source_vol.word_wrap = wrap
    txt_source_vol = txt_source_vol.paragraphs[0]
    txt_source_vol.font.name = font
    txt_source_vol.font.size = Pt(size)

    # creating runs (sections of text that can be individually formatted)
    run1 = txt_source_vol.add_run()
    run1.text = f'Across the time period, {company} was mentioned more often in '

    run2 = txt_source_vol.add_run()
    run2.text = f'{top_source_by_vol}'

    run3 = txt_source_vol.add_run()
    run3.text = ' than in any other publication. '
    
    # formatting above runs
    run2.font.italic = True

    print(f'\n\ntop_source_top_author = {top_source_top_author}')
    print(f'\n\ntop_source_top_author_vol = {top_source_top_author_vol}')

    if top_source_top_author != 'NA' and top_source_top_author_vol > 1: # top_source_top_author will be NA if unattributed
    
        # set volume figure as text rather than int digit if below 10 (must come after step above)
        if top_source_top_author != 'NA' and top_source_top_author_vol <10: top_source_top_author_vol = dict_low_numbers[top_source_top_author_vol]
        else: pass
    
        run3 = txt_source_vol.add_run()
        run3.text = f'{top_source_top_author} was the author most likely to mention {company_short} in this publication, with {top_source_top_author_vol} separate mentions.'

    else:
        """
        # Commented out. Function currently used for report over any period. Text below is therefore irrelevant for this purpose.
        if top_source_months_at_top < 12:
            
            if top_source_months_at_top <10: top_source_months_at_top = dict_low_numbers[top_source_months_at_top]
            else: pass
            
            run3 = txt_source_vol.add_run()
            run3.text = f'This was the most prominent source in {top_source_months_at_top} out of the 12 months of the year.'
            
        else:
            run3 = txt_source_vol.add_run()
            run3.text = f'This source was by far the leading generator of {company_short} mentions, being the most prominent publication in all 12 months of the year.'
        """

# %% add EOY text, top source by reach ----------------------------------------
    
def add_text_source_by_reach(slide, 
                           dataframe = df_temp,
            company = company, # variables set as themselves in text form are placeholders
            company_short = company_short, # these are to be replaced with the actual variables
            top_source_by_reach = top_source_by_reach,
            top_source_reach_text_long = top_source_reach_text_long,
            second_source_by_reach = second_source_by_reach,
            top_source_reach_pct_vs_2 = top_source_reach_pct_vs_2,
            top_source_by_vol = top_source_by_vol,
            top_source_reach_pct_vs_vol = top_source_reach_pct_vs_vol,
            left = Inches(1.5), 
            top = Inches(4.73), 
            width = Inches(5.54), 
            height = Inches(1.2), 
            font = 'Bahnschrift Light', 
            size = 16,
            wrap = True,
            align = PP_ALIGN.LEFT):
     
    txt_source_reach = (slide.shapes.add_textbox(left, top, width, height)).text_frame
    txt_source_reach.word_wrap = wrap
    txt_source_reach = txt_source_reach.paragraphs[0]
    txt_source_reach.font.name = font
    txt_source_reach.font.size = Pt(size)

    # creating runs (sections of text that can be individually formatted)
    if top_source_by_vol == top_source_by_reach:
        run1 = txt_source_reach.add_run()
        run1.text = f'{top_source_by_reach}'
        
        run2 = txt_source_reach.add_run()
        run2.text = f' was the also the widest-reaching publication. Its {top_source_reach_text_long} potential views was {top_source_reach_pct_vs_2}% higher than the number reached by the second widest reaching source, '
        
        run3 = txt_source_reach.add_run()
        run3.text = f'{second_source_by_reach}.'

        # format above runs
        run1.font.italic = True
        run3.font.italic = True
        
    else:
        run1 = txt_source_reach.add_run()
        run1.text = 'However, while '
        
        run2 = txt_source_reach.add_run()
        run2.text = f'{top_source_by_vol}'
        
        run3 = txt_source_reach.add_run()
        run3.text = f' was the most prolific publication, it was not the widest-reaching. The {top_source_reach_text_long} potential views reached by '
        
        run4 = txt_source_reach.add_run()
        run4.text = f'{top_source_by_reach}'
        
        if top_source_reach_pct_vs_vol == 'ERROR':
            run5 = txt_source_reach.add_run()
            run5.text = ' was the largest total audience reached by a single source.'
            
        else:
            run5 = txt_source_reach.add_run()
            run5.text = f' was {top_source_reach_pct_vs_vol} percent higher than the number reached by '
        
            run6 = txt_source_reach.add_run()
            run6.text = f'{top_source_by_vol}.'
            
            run6.font.italic = True
        
        # format above runs
        run2.font.italic = True
        run4.font.italic = True
        