# -*- coding: utf-8 -*-
"""
Created on Mon Nov 21 15:18:48 2022

@author: rober
"""

from utils.func_ppt import set_cell_border
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from utils.dicts_data_to_ppt import dict_media_colors

# %% table title height -------------------------------------------------------

# title height
def table_title_height(table,
                       title_height,
                       total_rows,
                       row_height):
                           
    table.rows[0].height = Inches (title_height)
    
    for row in range(1, total_rows):
        table.rows[row].height = Inches (row_height)
        
# %% format breakdown table text ----------------------------------------------
        
def format_breakdown_table_text(tbl):
    for col in range(0, 12):
        for row in range(0, 6):
            cell = tbl.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = 'Bahnschrift Light'
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(47, 85, 151)
            paragraph.font.bold = False
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for row in range(1, 6):
            cell = tbl.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        for row in range(3, 6):
            cell = tbl.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
        for row in range(4, 6):
            cell = tbl.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.italic = True
        for row in range(5, 6):
            cell = tbl.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(8)
            
# %% format table alt color ---------------------------------------------------            
            
# alt color columns
def format_table_alt_color(table,
                       rows,
                       cols,
                       col_color = RGBColor(242, 242, 242)):
    

    for row in range(0, rows):
            for col in range(0, cols):
                    table.cell(row, col).fill.background()
                    
    # set every other column to a light grey fill using third range argument
    for row in range(0, rows):
            for col in range(0, cols, 2):
                    table.cell(row, col).fill.solid()
                    table.cell(row, col).fill.fore_color.rgb = col_color

# %% format media types table -------------------------------------------------
                    
def format_media_types_table(table = 'default table', 
                             list_datapoints = [1,2,3,4,5,6],
                             ):
    # set borders to no borders
    for row in range(0, len(list_datapoints) + 1):
            for col in range(0, 2):
                set_cell_border(table.cell(row, col), "000000", '1')
    # set all table cells to transparent
    for row in range(0, len(list_datapoints) + 1):
            for col in range(0, 2):
                table.cell(row, col).fill.background()
    # set body text font, size, alignment and color
    for col in range(0, 2):
        for row in range(0, len(list_datapoints) + 1):
            cell = table.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = 'Bierstadt Display'
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.font.bold = False
            paragraph.alignment = PP_ALIGN.RIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for col in range(1, 2):
        for row in range(0, len(list_datapoints) + 1):
            cell = table.rows[row].cells[col]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(122, 122, 122)
            paragraph.alignment = PP_ALIGN.LEFT

# %% format media table text color --------------------------------------------

def format_media_text_color(table,
                       col_range_from, col_range_to):                    
    for col in range(col_range_from, col_range_to):
        try:
            for row in range(1, 2):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_1'])
        except:
            pass
        
        try:    
            for row in range(2, 3):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_2'])
        except:
            pass    
            
        try:    
            for row in range(3, 4):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_3'])
        except:
            pass    
            
        try:    
            for row in range(4, 5):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_4'])
        except:
            pass    
        
        try:    
            for row in range(5, 6):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_5'])
        except:
            pass
               
        try:    
            for row in range(6, 7):
                cell = table.rows[row].cells[col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor.from_string(dict_media_colors['media_other'])
        except:
            pass  
        