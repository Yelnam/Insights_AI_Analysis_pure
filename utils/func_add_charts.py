# -*- coding: utf-8 -*-
"""
Created on Mon Nov 21 11:19:44 2022

@author: rober
"""

# imports user-created PPT functions
from utils.func_ppt import set_reverse_categories
from utils.func_format_charts import format_y_axis, format_y_axis_bar, apply_data_labels
from utils.dicts_data_to_ppt import dict_media_colors

from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_LABEL_POSITION


# %% add_trend_chart - use for both vol and reach -----------------------------
def add_trend_chart(slide,
                    category_name = 'Volume',
                    data_series = [1,2,3],
                    left = Inches(4.32),
                    top = Inches(0.78),
                    height = Inches(8.35),
                    width = Inches(0.95),
                    chart_line_color = RGBColor(0,0,0)
                    ):
    
    
    # set categories and series data
    chart_trend_data = CategoryChartData()
    chart_trend_data.categories = [category_name]
    chart_trend_data.add_series(category_name, (data_series))
    
    # add chart to slide
    left, top, width, height = left, top, width, height
    trend_chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, 
                                    left, top, width, height, 
                                    chart_trend_data).chart
    
    # format y axis
    # function formats axis AND returns max y axis value to variable

    chart_max_y_for_label = format_y_axis(trend_chart, # chart
                data_series = data_series, # data series  
                visible = False,
                major_unit = 0, # will be auto-calculated if left at 0
                major_grid = True,
                major_grid_color = RGBColor(203, 203, 203),
                minor_grid = False,
                line_width = Inches(0),
                line_color = RGBColor(0,0,0),
                axis_max = 0) # will be auto-calculated if left at 0
        
    
    # format x axis
    xAxis = trend_chart.category_axis
    xAxis.visible = False
    
    # format data labels
    
    # apply different number formats depending on magnitude of data values
    if max(data_series) >= 1000000000:
        
        apply_data_labels(plot = trend_chart.plots[0],
                      font = 'Bahnschrift Light ',
                      fontsize = 10,
                      color = RGBColor(89,89,89),
                      numformat = "[<1000000000] 0.0,,\"m\";[>1000000000] 0.0,,,\"bn\"")
        
    elif max(data_series) >= 1000000:
        
        apply_data_labels(plot = trend_chart.plots[0],
                      font = 'Bahnschrift Light ',
                      fontsize = 10,
                      color = RGBColor(89,89,89),
                      numformat = "[<1000000] 0.0,k;[>1000000] 0.0,,\"m\"")
        
    elif max(data_series) >= 1000:
        
        apply_data_labels(plot = trend_chart.plots[0],
                      font = 'Bahnschrift Light ',
                      fontsize = 10,
                      color = RGBColor(89,89,89),
                      numformat = "[<1000] #;[>1000] #,##0")
        
    else:
        
        apply_data_labels(plot = trend_chart.plots[0],
                      font = 'Bahnschrift Light ',
                      fontsize = 10,
                      color = RGBColor(89,89,89),
                      numformat = "#")

            
    # set series colours
    
    # prevent PPT from using diff colours for same series
    trend_chart.plots[0].vary_by_categories = False
    
    # sets colours
    for series in trend_chart.series:
        series.format.line.color.rgb = chart_line_color
            
            
    # set series smoothness, legend, title
    trend_chart.series[0].smooth = True
    trend_chart.has_legend = False
    trend_chart.has_title = False
    
    return chart_max_y_for_label


# %% add_simple_pie -----------------------------------------------------------
def add_simple_pie(slide,
                    categories = ['Media Type 1', 'Media Type 2', 'Media Type 3'],
                    volumes = [1,2,3],
                    left = Inches(4.43),
                    top = Inches(3.45),
                    height = Inches(1.8),
                    width = Inches(1.8)
                    ):
    
   # defining chart data
   # creating object of chart
   media_types_chart_data = CategoryChartData()

   # adding categories to chart ----------------------------------

   media_types_chart_data.categories = categories

   # adding data series (category volumes) -----------------------

   # getting percentages so that pie chart labels are correctly formatted

   vols_pcts = [volumes[i] / sum(volumes)
                    for i in range(len(volumes))]

   media_types_chart_data.add_series('Volume',
                         (vols_pcts))
   
   simple_pie = slide.shapes.add_chart(
       XL_CHART_TYPE.PIE, left, top, width, height, media_types_chart_data
   ).chart


   simple_pie.has_title = False


   # legend removed to allow for standard chart size unaffected by legend
   #     colours indiicated in table text instead
   simple_pie.has_legend = False
   """
   simple_pie.legend.position = XL_LEGEND_POSITION.BOTTOM
   simple_pie.legend.include_in_layout = False
   simple_pie.legend.font.size = Pt(10)
   simple_pie.legend.font.color.rgb = RGBColor(122, 122, 122)
   """

   simple_pie.plots[0].has_data_labels = True
   data_labels = simple_pie.plots[0].data_labels
   data_labels.number_format = '0%'
   data_labels.position = XL_LABEL_POSITION.INSIDE_END
   data_labels.font.size = Pt(10)
   data_labels.font.name = 'Calibri'
   data_labels.font.color.rgb = RGBColor(255, 255, 255)


   # setting colours for categories
   points = simple_pie.plots[0].series[0].points

   # try/except if too few categories
       
   try:
       fill = points[0].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_1'])
   except:
       pass
       
       
   try:    
       fill = points[1].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_2'])
   except:
       pass    
       
       
   try:    
       fill = points[2].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_3'])
   except:
       pass    
       
       
   try:    
       fill = points[3].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_4'])
   except:
       pass    
       
       
   try:    
       fill = points[4].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_5'])
   except:
       pass


   try:    
       fill = points[5].format.fill
       fill.solid()
       fill.fore_color.rgb = RGBColor.from_string(dict_media_colors['media_other'])
   except:
       pass

# %% clustered_bar_noX --------------------------------------------------------
def clustered_bar_noX(slide, categories, volumes, left, top, width, height, chart_color = RGBColor(201, 179, 108)):
    # create chart -----------------------------------------------------
    
    # defining chart data
    # creating object of chart
    chart_data = CategoryChartData()
    
    # add categories and series to chart --------------------------
    chart_data.categories = categories
    chart_data.add_series("Volume", volumes)
    
    # adding clustered bar chart 
    chart_lead_cats = slide.shapes.add_chart(XL_CHART_TYPE.
                                                                BAR_CLUSTERED, 
                                               left, top, width, height, chart_data
                                               ).chart
    
    
    # setting chart physical attributes ---------------------------
    
    # set title and legend
    chart_lead_cats.has_title = False
    chart_lead_cats.has_legend = False
    
    # format x axis
    xAxis = chart_lead_cats.value_axis
    xAxis.visible = False
    xAxis.has_major_gridlines = False
    xAxis.has_minor_gridlines = False
    xAxis.minimum_scale = 0
    # axis max set 1.4 times max value to allow room for data labels
    # if no spokespeople present, set to 1
    try:
        xAxis.maximum_scale = (int(max(volumes)))*1.4
    except:
        xAxis.maximum_scale = 1
        
    """
    # commented out as it was producing a gap at top of y axis
    #  format y axis 
    format_y_axis_bar (chart = chart_lead_cats,
                       major_unit = 0,
                       major_grid_color = RGBColor(255,255,255),
                       major_grid = False,
                       minor_grid = False,
                       line_color = RGBColor(224, 224, 224),
                       label_size = 10.5,
                       label_color = RGBColor(122,122,122),
                       label_font = 'Bahnschrift Light',
                       axis_max = 5,
                       axis_min = 0)
    """

    yAxis = chart_lead_cats.category_axis
    yAxis.major_tick_mark = XL_TICK_MARK.NONE
    yAxis.format.line.width = Inches(0)
    yAxis.format.line.color.rgb = RGBColor(224, 224, 224)
    yAxis.tick_labels.font.name = 'Bahnschrift Light'
    yAxis.tick_labels.font.size = Pt(10.5)
    yAxis.tick_labels.font.color.rgb = RGBColor(122, 122, 122)

    # using function defined at start of PPT section of script to reverse Y axis
    set_reverse_categories(yAxis)

    # use imported PPT function to reverse Y axis
    set_reverse_categories(chart_lead_cats.category_axis) 
    
    """
    # setting series colours
    for series in chart_lead_cats.series:
        series.format.line.color.rgb = chart_color
    """
    
    # setting series colours
    colors = {}
    colors['Volume'] = chart_color
    
    for series in chart_lead_cats.series:
        if series.name in colors:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[series.name]
    
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = True
    
    # format labels - apply different number formats depending on magnitude of data values
    if max(volumes) >= 1000000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000000] 0.0,,\"m\";[>1000000000] 0.0,,,\"bn\"")
        
    elif max(volumes) >= 1000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000] 0.0,k;[>1000000] 0.0,,\"m\"")
        
    elif max(volumes) >= 1000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000] #;[>1000] #,##0")
        
    else:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "#")
        
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = False

# %% clustered_bar_noX_cats ---------------------------------------------------

# exactly as above, but uses as smaller text size for the y axis
def clustered_bar_noX_cats(slide, categories, volumes, left, top, width, height, chart_color = RGBColor(201, 179, 108)):
    # create chart -----------------------------------------------------
    
    # defining chart data
    # creating object of chart
    chart_data = CategoryChartData()
    
    # add categories and series to chart --------------------------
    chart_data.categories = categories
    chart_data.add_series("Volume", volumes)
    
    # adding clustered bar chart 
    chart_lead_cats = slide.shapes.add_chart(XL_CHART_TYPE.
                                                                BAR_CLUSTERED, 
                                               left, top, width, height, chart_data
                                               ).chart
    
    
    # setting chart physical attributes ---------------------------
    
    # set title and legend
    chart_lead_cats.has_title = False
    chart_lead_cats.has_legend = False
    
    # format x axis
    xAxis = chart_lead_cats.value_axis
    xAxis.visible = False
    xAxis.has_major_gridlines = False
    xAxis.has_minor_gridlines = False
    xAxis.minimum_scale = 0
    # axis max set 1.4 times max value to allow room for data labels
    # if no spokespeople present, set to 1
    try:
        xAxis.maximum_scale = (int(max(volumes)))*1.4
    except:
        xAxis.maximum_scale = 1
        
    # format y axis
    yAxis = chart_lead_cats.category_axis
    yAxis.major_tick_mark = XL_TICK_MARK.NONE
    yAxis.format.line.width = Inches(0)
    yAxis.format.line.color.rgb = RGBColor(224, 224, 224)
    yAxis.tick_labels.font.name = 'Bahnschrift Light'
    yAxis.tick_labels.font.size = Pt(7)
    yAxis.tick_labels.font.color.rgb = RGBColor(122, 122, 122)

    # using function defined at start of PPT section of script to reverse Y axis
    set_reverse_categories(yAxis)

    # use imported PPT function to reverse Y axis
    set_reverse_categories(chart_lead_cats.category_axis) 
    
    """
    # setting series colours
    for series in chart_lead_cats.series:
        series.format.line.color.rgb = chart_color
    """
    
    # setting series colours
    colors = {}
    colors['Volume'] = chart_color
    
    for series in chart_lead_cats.series:
        if series.name in colors:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[series.name]
    
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = True
    
    # format labels - apply different number formats depending on magnitude of data values
    if max(volumes) >= 1000000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000000] 0.0,,\"m\";[>1000000000] 0.0,,,\"bn\"")
        
    elif max(volumes) >= 1000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000] 0.0,k;[>1000000] 0.0,,\"m\"")
        
    elif max(volumes) >= 1000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000] #;[>1000] #,##0")
        
    else:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "#")
    
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = False
    
# %% clustered_bar_noX_sources ------------------------------------------------    
    
# exactly the same as clustered_bar_noX, but uses italics for the y axis 
def clustered_bar_noX_sources(slide, categories, volumes, left, top, width, height, chart_color = RGBColor(201, 179, 108)):
    # create chart -----------------------------------------------------
    
    # defining chart data
    # creating object of chart
    chart_data = CategoryChartData()
    
    # add categories and series to chart --------------------------
    chart_data.categories = categories
    chart_data.add_series("Volume", volumes)
    
    # adding clustered bar chart 
    chart_lead_cats = slide.shapes.add_chart(XL_CHART_TYPE.
                                                                BAR_CLUSTERED, 
                                               left, top, width, height, chart_data
                                               ).chart
    
    
    # setting chart physical attributes ---------------------------
    
    # set title and legend
    chart_lead_cats.has_title = False
    chart_lead_cats.has_legend = False
    
    # format x axis
    xAxis = chart_lead_cats.value_axis
    xAxis.visible = False
    xAxis.has_major_gridlines = False
    xAxis.has_minor_gridlines = False
    xAxis.minimum_scale = 0
    # axis max set 1.4 times max value to allow room for data labels
    # if no spokespeople present, set to 1
    try:
        xAxis.maximum_scale = (int(max(volumes)))*1.4
    except:
        xAxis.maximum_scale = 1
        
    """
    # commented out as it was producing a gap at top of y axis
    #  format y axis 
    format_y_axis_bar (chart = chart_lead_cats,
                       major_unit = 0,
                       major_grid_color = RGBColor(255,255,255),
                       major_grid = False,
                       minor_grid = False,
                       line_color = RGBColor(224, 224, 224),
                       label_size = 10.5,
                       label_color = RGBColor(122,122,122),
                       label_font = 'Bahnschrift Light',
                       axis_max = 5,
                       axis_min = 0)
    """

    yAxis = chart_lead_cats.category_axis
    yAxis.major_tick_mark = XL_TICK_MARK.NONE
    yAxis.format.line.width = Inches(0)
    yAxis.format.line.color.rgb = RGBColor(224, 224, 224)
    yAxis.tick_labels.font.name = 'Bahnschrift Light'
    yAxis.tick_labels.font.size = Pt(10.5)
    yAxis.tick_labels.font.color.rgb = RGBColor(122, 122, 122)
    yAxis.tick_labels.font.italic = True

    # using function defined at start of PPT section of script to reverse Y axis
    set_reverse_categories(yAxis)

    # use imported PPT function to reverse Y axis
    set_reverse_categories(chart_lead_cats.category_axis) 
    
    """
    # setting series colours
    for series in chart_lead_cats.series:
        series.format.line.color.rgb = chart_color
    """
    
    # setting series colours
    colors = {}
    colors['Volume'] = chart_color
    
    for series in chart_lead_cats.series:
        if series.name in colors:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[series.name]
    
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = True
    
    # format labels - apply different number formats depending on magnitude of data values
    if max(volumes) >= 1000000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000000] 0.0,,\"m\";[>1000000000] 0.0,,,\"bn\"")
        
    elif max(volumes) >= 1000000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000000] 0.0,k;[>1000000] 0.0,,\"m\"")
        
    elif max(volumes) >= 1000:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "[<1000] #;[>1000] #,##0")
        
    else:
        
        apply_data_labels(plot = chart_lead_cats.plots[0],
                          pos = XL_LABEL_POSITION.OUTSIDE_END,
                          font = 'Helvetica New ghd',
                          fontsize = 9,
                          color = RGBColor(127,127,127),
                      numformat = "#")
    
    # prevents PPT from using diff colours for same series
    chart_lead_cats.plots[0].vary_by_categories = False