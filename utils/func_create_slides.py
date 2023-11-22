import os 
from utils.dicts_data_to_ppt import dict_company_colors, dict_company_logos
from utils.func_add_charts import clustered_bar_noX
from utils.func_add_text import (add_text, add_text_EOY2022_title, add_text_EOY2022_header, 
                                 add_text_author_by_vol, add_text_source_by_vol)
from utils.lists_data_to_ppt import list_auth_format
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from itertools import zip_longest

# %% Create front page

def create_slide_front (ppt,
                        blank_slide_layout, 
                        company,
                        month_text):

    slide_front = ppt.slides.add_slide(blank_slide_layout) # adding a slide to the PPT file
    
    # ADD IMAGES ------------------------------------------
    
    # title slide image, background -------------------------------
    # set image dimensions
    left, top, width, height = Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    # add image
    image_path = os.path.join('.', 'Images', 'Backgrounds', 'on_front_page_bg.jpg')
    slide_front.shapes.add_picture(image_path, left, top, width, height)
    # title slide image, big logo ---------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    slide_front.shapes.add_picture(image_path, Inches(1), Inches(0.83), Inches(3.27), Inches(0.49)) # left, top, width, height
    # title slide image, icons ------------------------------------
    """
    # Commented out. Phone number no longer included on reports
    image_path = os.path.join('.', 'Images', 'Icons', 'phone.png')
    slide_front.shapes.add_picture(image_path, Inches(1.03), Inches(6.19), Inches(0.34), Inches(0.34))
    """

    image_path = os.path.join('.', 'Images', 'Icons', 'email.png')
    slide_front.shapes.add_picture(image_path, Inches(1.03), Inches(6.56), Inches(0.34), Inches(0.34))
    
    # ADD TEXT --------------------------------------------

    # title slide text, company name -------------------------------
    add_text(slide = slide_front, text = company, 
                left = Inches(0.91), top = Inches(3.12), 
                width = Inches(6.29), height = Inches(0.71), 
                font = 'Bahnschrift Semibold', size = 36, 
                color = RGBColor(000, 000, 000), wrap = True)
    # title slide text, report period ------------------------------
    add_text(slide = slide_front, text = f'Monthly Overview - {month_text}', 
                left = Inches(0.93), top = Inches(3.77), 
                width = Inches(6.29), height = Inches(0.71), 
                font = 'Bahnschrift Light', size = 24, 
                color = RGBColor(000, 000, 000), wrap = True)
    # title slide text, contact details --------------------------
    """
    # Commented out. Phone number no longer included on reports
    add_text(slide = slide_front, text = '+44 20 7264 4700', 
                left = Inches(1.36), top = Inches(6.25), 
                width = Inches(2.45), height = Inches(0.3), 
                font = 'Avenir Next LT Pro', size = 12, 
                color = RGBColor(000, 000, 000), wrap = True)
                """
    add_text(slide = slide_front, text = 'services@onclusive.com', 
                left = Inches(1.37), top = Inches(6.57), 
                width = Inches(2.45), height = Inches(0.3), 
                font = 'Avenir Next LT Pro', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)
    

# %% Create leading authors slide

def create_slide_authors(df,
                         ppt, 
                         blank_slide_layout, 
                         company, 
                         company_short,
                         company_color,
                         top_author_by_vol,
                         top_author_by_vol_surname, 
                         top_author_by_vol_source,
                         top_author_vol,
                         list_top_authors,
                         list_top_authors_vols
                         ):

    # %%%% LEADING AUTHORS - ADD SLIDE TO PPT FILE ---------------------------
    slide_authors = ppt.slides.add_slide(blank_slide_layout)
    slide_authors.background.fill.solid()
    slide_authors.background.fill.fore_color.rgb = RGBColor(204, 255, 184)
    # %%%% LEADING AUTHORS - ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide_authors.shapes.add_picture(image_path,
                        Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    image_path = os.path.join('.', 'Images', 'Logos', f'{company}.png')
    company_logo = (slide_authors.shapes.add_picture(image_path, 
                        left=Inches(12.32), top=Inches(0.23), height=Inches(0.98), width=Inches(0.7)))
    
    # %%%% LEADING AUTHORS - ADD TEXT ----------------------------------------
    add_text_EOY2022_title(slide = slide_authors, text = 'Leading Authors')
    # headers
    add_text_EOY2022_header(slide = slide_authors, 
                            text = 'Top Authors by Volume', left = Inches(7.57), top = Inches(2.27))

    # body text
    # author vol text
    add_text_author_by_vol(slide = slide_authors, dataframe = df,
                company = company, company_short = company_short, #
                top_author_by_vol = top_author_by_vol,
                top_author_by_vol_surname = top_author_by_vol_surname,
                top_author_by_vol_source = top_author_by_vol_source,
                top_author_vol = top_author_vol,
                left = Inches(1.5), top = Inches(3.22), 
                width = Inches(5.54), height = Inches(1.2))
    
    # %%%% LEADING AUTHORS - ADD CHARTS --------------------------------------
    chart_journo_vol = clustered_bar_noX(slide = slide_authors, #create a clustered bar chart with invisible X axis
                    categories = list_top_authors[0:5], volumes = list_top_authors_vols[0:5],
                    left = Inches(7.72), top = Inches(2.52),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create leading authors slide

def create_slide_sources(df,
                         ppt, 
                         blank_slide_layout, 
                         company, 
                         company_short,
                         company_color,
                         top_source_by_vol,
                         top_source_vol, 
                         top_source_top_author,
                         top_source_top_author_vol,
                         # top_source_vol_text_long,
                         list_top_sources_by_vol,
                         list_top_sources_vol
                         ):

    # %%%% LEADING SOURCES - ADD SLIDE TO PPT FILE ---------------------------
    slide_sources = ppt.slides.add_slide(blank_slide_layout)
    slide_sources.background.fill.solid()
    slide_sources.background.fill.fore_color.rgb = RGBColor(204, 255, 184)
    # %%%% LEADING SOURCES - ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide_sources.shapes.add_picture(image_path,
                        Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    image_path = os.path.join('.', 'Images', 'Logos', f'{company}.png')
    company_logo = (slide_sources.shapes.add_picture(image_path, 
                        left=Inches(12.32), top=Inches(0.23), height=Inches(0.98), width=Inches(0.7)))
    
    # %%%% LEADING SOURCES - ADD TEXT ----------------------------------------
    add_text_EOY2022_title(slide = slide_sources, text = 'Leading Sources')
    # headers
    add_text_EOY2022_header(slide = slide_sources, 
                            text = 'Top Sources by Volume', left = Inches(7.57), top = Inches(2.27))

    # body text
    # source vol text
    add_text_source_by_vol(slide = slide_sources, dataframe = df,
                company = company, company_short = company_short, #
                top_source_by_vol = top_source_by_vol,
                top_source_top_author = top_source_top_author,
                top_source_top_author_vol = top_source_top_author_vol,
                left = Inches(1.5), top = Inches(3.22), 
                width = Inches(5.54), height = Inches(1.2))
    
    # %%%% LEADING SOURCES - ADD CHARTS --------------------------------------
    chart_source_vol = clustered_bar_noX(slide = slide_sources, #create a clustered bar chart with invisible X axis
                    categories = list_top_sources_by_vol[0:5], volumes = list_top_sources_vol[0:5],
                    left = Inches(7.72), top = Inches(2.52),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    

# %% Create brand corporate/consumer slide
def create_slide_corp_con (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_corp_con,
                            dict_brand_corp_con_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(254, 246, 214)

    # remove line
    header_bg.line.fill.background()

    # remove shadow
    header_bg.shadow.inherit = False
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Corporate/Consumer Breakdown - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Corp/Cons Breakdown by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_corp_con_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')
    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories = ['Corporate', 'Consumer'], 
                    volumes = [
                                dict_brand_corp_con[brand].get('Corporate', 0),
                                dict_brand_corp_con[brand].get('Consumer', 0)
                            ],                            
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create brand sentiment slide
def create_slide_sentiment (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_sentiment,
                            dict_brand_sent_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(212, 252, 225)

    # remove line
    header_bg.line.fill.background()

    # remove shadow
    header_bg.shadow.inherit = False

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))
    
    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Sentiment - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Sentiment by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_sent_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')
    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories = ['Positive', 'Neutral', 'Negative'], 
                    volumes = [
                                dict_brand_sentiment[brand].get('Positive', 0),
                                dict_brand_sentiment[brand].get('Neutral', 0),
                                dict_brand_sentiment[brand].get('Negative', 0)
                            ],
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create brand prominence slide
def create_slide_prominence (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_prominence,
                            dict_brand_prom_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(219, 229, 241)

    # remove line
    header_bg.line.fill.background()

    # remove shadow
    header_bg.shadow.inherit = False
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Prominence - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Prominence by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_prom_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')
    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories = ['Headline', 'First paragraph', 'Top half of article', 'Bottom half of article'], 
                    volumes = [
                                dict_brand_prominence[brand].get('Headline', 0),
                                dict_brand_prominence[brand].get('First paragraph', 0),
                                dict_brand_prominence[brand].get('Top half of article', 0),
                                dict_brand_prominence[brand].get('Bottom half of article', 0)
                            ],
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    

# %% Create brand topics slide
def create_slide_topics (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_topics,
                            dict_brand_topics_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(238, 238, 238)

    header_bg.line.fill.background()    # Remove line
    header_bg.shadow.inherit = False    # Remove shadow
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Topics - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Topics by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_topics_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')

    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories = ['Financial Performance', 'Regulation', 'Product launch', 'Offers or Deals'], 
                    volumes = [
                                dict_brand_topics[brand].get('Financial Performance', 0),
                                dict_brand_topics[brand].get('Regulation', 0),
                                dict_brand_topics[brand].get('Product launch', 0),
                                dict_brand_topics[brand].get('Offers or Deals', 0)
                            ],
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create brand pos values slide
def create_slide_pos_values (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_pos_values,
                            dict_brand_pos_values_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(219, 255, 205)

    header_bg.line.fill.background()    # Remove line
    header_bg.shadow.inherit = False    # Remove shadow
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Positive Brand Values - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Pos. BVs by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_pos_values_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')

    categories =  ['Design', 'Innovation', 'Inclusivity', 'Sustainability', 'Transparency/Trust', 'Value']
    volumes = [
                                dict_brand_pos_values[brand].get(categories[0], 0),
                                dict_brand_pos_values[brand].get(categories[1], 0),
                                dict_brand_pos_values[brand].get(categories[2], 0),
                                dict_brand_pos_values[brand].get(categories[3], 0),
                                dict_brand_pos_values[brand].get(categories[4], 0),
                                dict_brand_pos_values[brand].get(categories[5], 0)
                            ]
    
    print(dict_brand_pos_values[brand])
    print(dict_brand_pos_values[brand].get('Design', 0))
    print(dict_brand_pos_values[brand].get('Innovation', 0))
    print(dict_brand_pos_values[brand].get('Inclusivity', 0))
    print(dict_brand_pos_values[brand].get('Sustainability', 0))
    print(dict_brand_pos_values[brand].get('Transparency/Trust', 0))
    print(dict_brand_pos_values[brand].get('Value', 0))
    
    
    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories =  categories, 
                    volumes = volumes,
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create brand neg values slide
def create_slide_neg_values (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_neg_values,
                            dict_brand_neg_values_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(253, 230, 211)

    header_bg.line.fill.background()    # Remove line
    header_bg.shadow.inherit = False    # Remove shadow
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Negative Brand Values - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Neg. BVs by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_neg_values_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')

    categories =  ['Design', 'Innovation', 'Inclusivity', 'Sustainability', 'Transparency/Trust', 'Value']
    volumes = [
                                dict_brand_neg_values[brand].get(categories[0], 0),
                                dict_brand_neg_values[brand].get(categories[1], 0),
                                dict_brand_neg_values[brand].get(categories[2], 0),
                                dict_brand_neg_values[brand].get(categories[3], 0),
                                dict_brand_neg_values[brand].get(categories[4], 0),
                                dict_brand_neg_values[brand].get(categories[5], 0)
                            ]
    
    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories =  categories, 
                    volumes = volumes,
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)
    
# %% Create brand spokespeople slide
def create_slide_spokespeople (ppt,
                            blank_slide_layout, 
                            brand,
                            dict_brand_spokespeople,
                            dict_brand_spokespeople_text,
                            company_color_yn,
                            company_logo_yn) :
    
    slide = ppt.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ADD HEADER BACKGROUND SHAPE ---------------------

    shape = MSO_SHAPE.RECTANGLE
    left = 0
    top = 0
    width = Inches(13.33)
    height = Inches(1.21)

    header_bg = slide.shapes.add_shape(shape, left, top, width, height)

    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(226, 230, 254)

    header_bg.line.fill.background()    # Remove line
    header_bg.shadow.inherit = False    # Remove shadow
    

    # ADD IMAGES --------------------------------------
    image_path = os.path.join('.', 'Images', 'Logos', 'onclusive_logo_black.png')
    onclusive_logo = (slide.shapes.add_picture(image_path, Inches(0.41), Inches(6.94), Inches(1.81), Inches(0.27)))
    
    if company_logo_yn == 'co_logo_Y':
        image_path = os.path.join('.', 'Images', 'Logos', f'{brand}.png')
        company_logo = (slide.shapes.add_picture(image_path, 
                                                left = Inches(dict_company_logos[f'{brand}'][0]), top = Inches(dict_company_logos[f'{brand}'][1]),
                                                width = Inches(dict_company_logos[f'{brand}'][3]), height = Inches(dict_company_logos[f'{brand}'][2])))

    # ADD TEXT ----------------------------------------

    # title
    title = add_text_EOY2022_title(slide = slide, text = f'Spokespeople - {brand}', left = Inches(0.42), top = Inches(0.34)) # slide title

    # headers
    header = add_text_EOY2022_header(slide = slide, 
                text = f'Spokespeople by Volume - {brand}', left = Inches(7.86), top = Inches(2.05))

    # body text
    add_text(slide = slide, text = dict_brand_spokespeople_text[brand], 
                left = Inches(0.91), top = Inches(2.06), 
                width = Inches(5.9), height = Inches(4), 
                font = 'Bahnschrift Light', size = 14, 
                color = RGBColor(000, 000, 000), wrap = True)

    # ADD CHARTS --------------------------------------

    # adjust company_color for brand
    if company_color_yn == 'co_color_Y':
        company_color = RGBColor.from_string(dict_company_colors[f'{brand}'])
    else:
        company_color = RGBColor.from_string('B788FA')
    
    spokespeople_counts = dict_brand_spokespeople[brand]
    top_spokespeople = spokespeople_counts.most_common(4)

    # Extend top_spokespeople to always contain 4 elements
    extended_top_spokespeople = list(zip_longest(top_spokespeople, range(4)))

    categories = [spokesperson if spokesperson is not None else '' for spokesperson, _ in extended_top_spokespeople]

    volumes = []
    for i in range(4):
        try:
            volumes.append(spokespeople_counts.get(categories[i], 0))
        except IndexError:
            volumes.append(0)

    chart = clustered_bar_noX(slide = slide, #create a clustered bar chart with invisible X axis
                    categories = categories, 
                    volumes = volumes,
                    left = Inches(8.04), top = Inches(2.5),
                    width = Inches(3.69), height = Inches(2.71),
                    chart_color = company_color)